"""知识库文档加载器的公开行为契约。"""

import base64
from types import SimpleNamespace
from unittest.mock import patch

import fitz
import pandas as pd
import pydantic.root_model  # noqa: F401
import pytest
from docx import Document as WordDocument
from langchain_core.documents import Document
from pptx import Presentation
from pptx.util import Inches

from apps.opspilot.metis.llm.loader import pdf_loader, website_loader
from apps.opspilot.metis.llm.loader.doc_loader import DocLoader
from apps.opspilot.metis.llm.loader.excel_loader import ExcelLoader
from apps.opspilot.metis.llm.loader.pdf_loader import PDFLoader
from apps.opspilot.metis.llm.loader.ppt_loader import PPTLoader
from apps.opspilot.metis.llm.loader.website_loader import WebSiteLoader


pytestmark = pytest.mark.unit

PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


class OCRFromFile:
    def __init__(self, result="diagram text"):
        self.result = result
        self.paths = []

    def predict(self, path):
        self.paths.append(path)
        return self.result


class OCRFromBase64:
    def __init__(self, result="image text"):
        self.result = result
        self.images = []

    def predict_from_base64(self, image):
        self.images.append(image)
        return self.result


def write_png(path):
    path.write_bytes(PNG_1X1)


def test_excel_loader_full_mode_preserves_column_names_and_cell_values(tmp_path):
    workbook = tmp_path / "inventory.xlsx"
    with pd.ExcelWriter(workbook, engine="openpyxl") as writer:
        pd.DataFrame(
            [
                {"name": "database-01", "count": 3},
                {"name": "cache-01", "count": 5},
            ]
        ).to_excel(writer, sheet_name="assets", index=False)

    docs = ExcelLoader(str(workbook), mode="full").load()

    assert len(docs) == 1
    assert docs[0].metadata == {
        "format": "table",
        "sheets": ["assets"],
        "source": str(workbook),
    }
    assert "name\tcount" in docs[0].page_content
    assert "database-01\t3" in docs[0].page_content


def test_excel_loader_supports_structured_row_and_per_sheet_modes(tmp_path):
    workbook = tmp_path / "metrics.xlsx"
    with pd.ExcelWriter(workbook, engine="openpyxl") as writer:
        pd.DataFrame([{"host": "db-01", "cpu": 91}]).to_excel(
            writer,
            sheet_name="critical",
            index=False,
        )
        pd.DataFrame([{"host": "web-01", "cpu": 12}]).to_excel(
            writer,
            sheet_name="healthy",
            index=False,
        )

    rows = ExcelLoader(str(workbook), mode="excel_header_row_parse").load()
    sheets = ExcelLoader(str(workbook), mode="excel_full_content_parse").load()

    assert [doc.metadata["sheet"] for doc in rows] == ["critical", "healthy"]
    assert rows[0].page_content == "critical  host: db-01  critical  cpu: 91"
    assert [doc.metadata["sheet"] for doc in sheets] == ["critical", "healthy"]
    assert sheets[1].metadata["source"] == str(workbook)
    assert "web-01\t12" in sheets[1].page_content


def test_excel_loader_rejects_unknown_mode():
    with pytest.raises(ValueError, match="Unsupported mode"):
        ExcelLoader("unused.xlsx", mode="rows").load()


def test_doc_loader_reads_paragraphs_tables_and_images(tmp_path):
    source = tmp_path / "runbook.docx"
    image = tmp_path / "diagram.png"
    write_png(image)
    document = WordDocument()
    document.add_heading("Database Runbook", level=1)
    document.add_paragraph("Restart only after draining traffic.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "service"
    table.cell(0, 1).text = "owner"
    table.cell(1, 0).text = "postgres"
    table.cell(1, 1).text = "dba"
    document.add_picture(str(image))
    document.save(source)
    ocr = OCRFromFile()

    docs = DocLoader(str(source), ocr=ocr, mode="paragraph").load()

    assert docs[0].page_content == (
        "Database Runbook\nRestart only after draining traffic."
    )
    assert docs[1].metadata == {"format": "table"}
    assert "| postgres | dba |" in docs[1].page_content
    assert docs[2].metadata["format"] == "image"
    assert docs[2].page_content == "diagram text"
    assert len(ocr.paths) == 1
    assert not __import__("os").path.exists(ocr.paths[0])


def test_doc_loader_full_mode_and_invalid_document_contract(tmp_path):
    source = tmp_path / "notes.docx"
    document = WordDocument()
    document.add_paragraph("first")
    document.add_paragraph("second")
    document.save(source)

    docs = DocLoader(str(source), ocr=None, mode="full").load()
    invalid_mode = DocLoader(str(source), ocr=None, mode="unsupported").load()
    missing = DocLoader(str(tmp_path / "missing.docx"), ocr=None).load()

    assert [doc.page_content for doc in docs] == ["firstsecond"]
    assert invalid_mode == []
    assert missing == []


def test_ppt_loader_reads_page_text_table_and_image(tmp_path):
    source = tmp_path / "operations.pptx"
    image = tmp_path / "architecture.png"
    write_png(image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(4),
        Inches(1),
    )
    textbox.text_frame.text = "Incident response"
    table = slide.shapes.add_table(
        1,
        2,
        Inches(0.5),
        Inches(2),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = "severity"
    table.cell(0, 1).text = "P1"
    slide.shapes.add_picture(str(image), Inches(5), Inches(0.5))
    presentation.save(source)
    ocr = OCRFromBase64()

    docs = PPTLoader(str(source), ocr=ocr, load_mode="page").load()

    formats = [doc.metadata.get("format") for doc in docs]
    assert formats == ["table", "image", None]
    assert docs[0].page_content == "severity\nP1"
    assert docs[1].page_content == "image text"
    assert docs[1].metadata["page"] == 1
    assert docs[2].page_content == "Incident response"
    assert len(ocr.images) == 1


def test_ppt_loader_full_mode_combines_slide_text(tmp_path):
    source = tmp_path / "summary.pptx"
    presentation = Presentation()
    for text in ("Capacity", "Availability"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(4),
            Inches(1),
        )
        textbox.text_frame.text = text
    presentation.save(source)

    docs = PPTLoader(str(source), ocr=None, load_mode="full").load()

    assert len(docs) == 1
    assert docs[0].page_content == "Capacity\nAvailability"


def test_pdf_loader_reads_real_pages_and_ocr_image(tmp_path):
    source = tmp_path / "manual.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Maintenance window")
    page.insert_image(fitz.Rect(72, 100, 92, 120), stream=PNG_1X1)
    pdf.save(source)
    pdf.close()
    ocr = OCRFromFile("topology")

    with patch.object(pdf_loader, "read_pdf", side_effect=[[], []]):
        docs = PDFLoader(str(source), ocr=ocr, load_mode="page").load()

    assert any(
        doc.metadata == {"format": "text", "page": 1}
        and "Maintenance window" in doc.page_content
        for doc in docs
    )
    image_doc = next(doc for doc in docs if doc.metadata["format"] == "image")
    assert image_doc.page_content == "topology"
    assert image_doc.metadata["page"] == 1
    assert len(ocr.paths) == 1
    assert not __import__("os").path.exists(ocr.paths[0])


def test_pdf_loader_full_mode_ignores_overlapping_table_text(tmp_path):
    source = tmp_path / "report.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "outside")
    pdf.save(source)
    pdf.close()

    with patch.object(pdf_loader, "read_pdf", side_effect=[[], []]):
        docs = PDFLoader(str(source), ocr=None, load_mode="full").load()

    assert [doc.page_content for doc in docs] == ["outside"]


def test_website_loader_transforms_html_and_ocr_images():
    web_docs = [
        Document(
            "<html><body><h1>Status</h1><p>Healthy</p>"
            '<img src="/diagram.png" alt="Topology"></body></html>',
            metadata={"source": "https://93.184.216.34/docs/index.html"},
        )
    ]

    class ExternalCrawler:
        def __init__(self, url, **kwargs):
            self.url = url
            self.kwargs = kwargs

        def load(self):
            return web_docs

    ocr = OCRFromBase64("primary -> replica")
    response = SimpleNamespace(
        status_code=200,
        headers={"content-type": "image/png"},
        content=PNG_1X1,
    )
    with (
        patch.object(website_loader, "RecursiveUrlLoader", ExternalCrawler),
        patch.object(website_loader, "safe_get", return_value=response),
    ):
        docs = WebSiteLoader(
            "https://93.184.216.34/docs/index.html",
            max_depth=2,
            ocr=ocr,
        ).load()

    assert docs[0].page_content == "Healthy"
    assert docs[1].metadata["format"] == "image"
    assert docs[1].metadata["page"] == 1
    assert "图片描述(alt): Topology" in docs[1].page_content
    assert "primary -> replica" in docs[1].page_content
    assert len(ocr.images) == 1
